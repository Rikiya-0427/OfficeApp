# prs = Presentation()
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

# title.text = "Hello, World!"
# subtitle.text = "python-pptx was here!"

# prs.save('test.pptx')

from pptx import Presentation

prs = Presentation()
title_slide_layout = prs.slide_layouts[1]  # スライドレイアウトのインデックスを変更
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Hello, OpenAI!"  # タイトルを変更
subtitle.text = "TAPIOKA was here!"  # サブタイトルを変更

prs.save('test2.pptx')

# 作成したパワーポイントファイル内の文字数をカウントして表示する
presentation = Presentation('test2.pptx')
text_count = 0
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_count += len(run.text.strip())

print("文字数:", text_count)
